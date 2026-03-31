"""Document generation engine for populating Word and PowerPoint templates.

Reads JSON template specs and fills templates with computed fund data,
charts, and metrics tables using python-docx and python-pptx.
"""

from __future__ import annotations

import io
import json
import re
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from docx import Document
from docx.shared import Emu
from pptx import Presentation

from fofproject.batch import plot_cumulative_returns
from fofproject.fund import Fund, save_dir
from fofproject.paths import TEMPLATE_DIR
from fofproject.utils import parse_month

# Chinese month names for date formatting
CN_MONTH_NAMES = {
    1: "一月", 2: "二月", 3: "三月", 4: "四月",
    5: "五月", 6: "六月", 7: "七月", 8: "八月",
    9: "九月", 10: "十月", 11: "十一月", 12: "十二月",
}

EN_MONTH_NAMES = {
    1: "January", 2: "February", 3: "March", 4: "April",
    5: "May", 6: "June", 7: "July", 8: "August",
    9: "September", 10: "October", 11: "November", 12: "December",
}


class TemplateEngine:
    """Core engine that resolves specs and populates templates."""

    def __init__(
        self,
        funds: Dict[str, Fund],
        report_month: str,
        spec: dict,
    ):
        self.funds = funds
        self.report_month_str = report_month
        self.report_month = parse_month(report_month)
        self.spec = spec
        self.params = spec.get("parameters", {})
        self.language = self.params.get("language", "en")

    # ------------------------------------------------------------------ #
    #  Expression resolution
    # ------------------------------------------------------------------ #

    def resolve(self, expr) -> Any:
        """Resolve template expressions like {report_month}, {fund.inception_date}."""
        if expr is None:
            return None
        if not isinstance(expr, str):
            return expr
        if not expr.startswith("{") or not expr.endswith("}"):
            return expr

        key = expr[1:-1]

        if key == "report_month":
            return self.report_month
        if key == "year_start":
            return datetime(self.report_month.year, 1, 1)
        if key == "language":
            return self.language

        # {fund.inception_date} — needs fund context from element
        if key.startswith("fund."):
            return key  # handled in compute_value with fund context

        # {benchmarks.xxx} — lookup from parameters
        if key.startswith("benchmarks."):
            parts = key.split(".", 1)[1]
            return self.params["benchmarks"].get(parts)

        # {cumulative_return_chart_funds}
        if key in self.params:
            return self.params[key]

        return expr

    def resolve_args(self, args: dict, fund: Optional[Fund] = None) -> dict:
        """Resolve all expressions in an args dict."""
        resolved = {}
        for k, v in args.items():
            if isinstance(v, str) and v.startswith("{") and v.endswith("}"):
                key = v[1:-1]
                if key == "fund.inception_date" and fund:
                    resolved[k] = fund.inception_date
                elif key == "fund.latest_date" and fund:
                    resolved[k] = fund.latest_date
                else:
                    resolved[k] = self.resolve(v)
            else:
                resolved[k] = v
        return resolved

    # ------------------------------------------------------------------ #
    #  Value computation
    # ------------------------------------------------------------------ #

    def compute_value(self, data_source: dict) -> Any:
        """Call the appropriate Fund method and return the raw value."""
        method = data_source.get("method")
        if method in ("format_report_date", "manual_image"):
            return self._format_report_date(data_source)

        fund_name = data_source.get("fund")
        fund = self.funds.get(fund_name) if fund_name else None
        args = self.resolve_args(data_source.get("args", {}), fund)

        if method == "annualized_return":
            return fund.annualized_return(args["start_month"], args["end_month"])
        elif method == "cumulative_return":
            return fund.cumulative_return(args["start_month"], args["end_month"])
        elif method == "volatility":
            return fund.volatility(args["start_month"], args["end_month"])
        elif method == "sharpe_ratio":
            return fund.sharpe_ratio(args["start_month"], args["end_month"])
        elif method == "sortino_ratio":
            return fund.sortino_ratio(args["start_month"], args["end_month"])
        elif method == "max_drawdown":
            return fund.max_drawdown(args["start_month"], args["end_month"])
        elif method == "positive_months":
            return fund.positive_months(args["start_month"], args["end_month"])
        elif method == "return_in_positive_months":
            return fund.return_in_positive_months(args["start_month"], args["end_month"])
        elif method == "return_in_negative_months":
            return fund.return_in_negative_months(args["start_month"], args["end_month"])
        elif method == "beta_to":
            benchmark_name = data_source.get("benchmark_fund")
            benchmark = self.funds[benchmark_name]
            return fund.beta_to(benchmark, args["start_month"], args["end_month"])
        elif method == "get_monthly_return":
            return fund.get_monthly_return(args["year"], args["month"])
        else:
            raise ValueError(f"Unknown method: {method}")

    def _format_report_date(self, data_source: dict) -> str:
        """Format the report month as a date string."""
        fmt = data_source.get("args", {}).get("format", "MMMM YYYY")
        m = self.report_month.month
        y = self.report_month.year

        if fmt == "MMMM YYYY":
            return f"{EN_MONTH_NAMES[m]} {y}"
        elif fmt == "MMMM":
            return EN_MONTH_NAMES[m]
        elif fmt == "YYYY":
            return str(y)
        elif fmt == "YYYY年M月_cn":
            return f"{y}年{CN_MONTH_NAMES[m]}"
        else:
            return f"{EN_MONTH_NAMES[m]} {y}"

    # ------------------------------------------------------------------ #
    #  Value formatting
    # ------------------------------------------------------------------ #

    def format_value(self, value: Any, fmt: dict) -> str:
        """Format a numeric value according to spec rules."""
        if value is None:
            return fmt.get("empty_value", "-")

        fmt_type = fmt.get("type", "decimal")
        decimals = fmt.get("decimals", 2)
        multiply = fmt.get("multiply_100", False)
        negate = fmt.get("negate", False)
        suffix = fmt.get("suffix", "")

        if multiply:
            value = value * 100
        if negate:
            value = -value

        if fmt_type in ("percent", "percent_space"):
            formatted = f"{value:.{decimals}f}%"
            if fmt_type == "percent_space":
                formatted = f"{value:.{decimals}f} %"
        elif fmt_type == "percent_no_symbol":
            formatted = f"{value:.{decimals}f}"
            formatted += suffix
        elif fmt_type == "percent_signed":
            sign = "+" if value > 0 else ""
            formatted = f"{sign}{value:.{decimals}f}%"
            formatted += suffix.replace("%", "")  # avoid double %
        elif fmt_type == "decimal":
            formatted = f"{value:.{decimals}f}"
            formatted += suffix
        else:
            formatted = str(value)

        return formatted

    # ------------------------------------------------------------------ #
    #  Monthly returns grid builder
    # ------------------------------------------------------------------ #

    def build_monthly_returns_grid(
        self,
        fund: Fund,
        benchmark_fund: Fund,
        benchmark_label: str,
        end_month: datetime,
        num_years: Optional[int],
        include_ytd: bool,
        include_inception: bool,
        fmt: dict,
    ) -> List[List[str]]:
        """Build a 2D grid of monthly returns for populating table ranges.

        Returns rows in pairs: [year_fund_row, benchmark_row, year_fund_row, ...]
        Each row has 15 columns: [label, Jan..Dec, YTD, Since Inception]
        """
        end_year = end_month.year

        # Determine start year
        inception_year = fund.inception_date.year
        if num_years:
            start_year = end_year - num_years + 1
            start_year = max(start_year, inception_year)
        else:
            start_year = inception_year

        rows = []
        for year in range(end_year, start_year - 1, -1):
            # Fund row
            fund_row = [str(year)]
            for month in range(1, 13):
                # Check if this month is in the fund's data range
                month_dt = datetime(year, month, 1)
                if month_dt > end_month:
                    fund_row.append("")
                elif month_dt < fund.inception_date:
                    fund_row.append(fmt.get("empty_value", "-"))
                else:
                    val = fund.get_monthly_return(year, month)
                    if val is not None:
                        fund_row.append(self.format_value(val, fmt))
                    else:
                        fund_row.append(fmt.get("empty_value", "-"))

            # YTD column
            if include_ytd:
                jan = datetime(year, 1, 1)
                ytd_end = end_month if year == end_year else datetime(year, 12, 1)
                # Only compute if fund has data in this year
                if jan >= fund.inception_date:
                    ytd = fund.cumulative_return(jan, ytd_end)
                    fund_row.append(self.format_value(ytd, fmt))
                else:
                    # Partial year — compute from inception month
                    ytd = fund.cumulative_return(fund.inception_date, ytd_end)
                    fund_row.append(self.format_value(ytd, fmt))

            # Since inception column
            if include_inception:
                ytd_end = end_month if year == end_year else datetime(year, 12, 1)
                inception_val = fund.cumulative_return(fund.inception_date, ytd_end)
                fund_row.append(self.format_value(inception_val, fmt))

            rows.append(fund_row)

            # Benchmark row
            bench_row = [benchmark_label]
            for month in range(1, 13):
                month_dt = datetime(year, month, 1)
                if month_dt > end_month:
                    bench_row.append("")
                else:
                    val = benchmark_fund.get_monthly_return(year, month)
                    if val is not None:
                        bench_row.append(self.format_value(val, fmt))
                    else:
                        bench_row.append(fmt.get("empty_value", "-"))

            # Benchmark YTD
            if include_ytd:
                jan = datetime(year, 1, 1)
                ytd_end = end_month if year == end_year else datetime(year, 12, 1)
                try:
                    bench_ytd = benchmark_fund.cumulative_return(jan, ytd_end)
                    bench_row.append(self.format_value(bench_ytd, fmt))
                except Exception:
                    bench_row.append(fmt.get("empty_value", "-"))

            # Benchmark since inception — leave empty (not meaningful for benchmark)
            if include_inception:
                bench_row.append("")

            rows.append(bench_row)

        return rows

    # ------------------------------------------------------------------ #
    #  Calendar year table builder (for PPTX slide 13)
    # ------------------------------------------------------------------ #

    def build_calendar_year_table(
        self,
        funds_map: Dict[str, str],
        start_year: int,
        end_month: datetime,
        end_year_offset: int,
        fmt: dict,
    ) -> List[List[str]]:
        """Build calendar year performance table data.

        Returns a 2D list where each row is [Year, Fund1%, Fund2%, ...].
        """
        end_year = end_month.year + end_year_offset
        rows = []
        for year in range(start_year, end_year + 1):
            row = [str(year)]
            jan = datetime(year, 1, 1)
            dec = datetime(year, 12, 1)
            for label, fund_name in funds_map.items():
                if label == "Year":
                    continue
                fund = self.funds.get(fund_name)
                if fund:
                    try:
                        val = fund.cumulative_return(jan, dec)
                        row.append(self.format_value(val, fmt))
                    except Exception:
                        row.append("-")
                else:
                    row.append("-")
            rows.append(row)
        return rows

    # ------------------------------------------------------------------ #
    #  Image generation helpers
    # ------------------------------------------------------------------ #

    def generate_chart_image(self, data_source: dict) -> Optional[bytes]:
        """Generate a chart/table image and return as bytes."""
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        method = data_source.get("method")
        args = data_source.get("args", {})

        if method == "plot_cumulative_returns":
            fund_names = data_source.get("funds", [])
            # Resolve if it's an expression
            if isinstance(fund_names, str):
                fund_names = self.resolve(fund_names)
            subset = {n: self.funds[n] for n in fund_names if n in self.funds}
            resolved_args = {}
            for k, v in args.items():
                resolved_args[k] = self.resolve(v) if isinstance(v, str) else v
            # Convert end_month to string for batch.py
            em = resolved_args.get("end_month")
            if isinstance(em, datetime):
                resolved_args["end_month"] = em.strftime("%Y-%m")
            resolved_args["save"] = False
            fig = plot_cumulative_returns(funds=subset, **resolved_args)
            # Plotly figure — convert to PNG bytes using the figure's own
            # layout dimensions (set by aspect_lock in plot_cumulative_returns).
            return fig.to_image(format="png")

        elif method == "export_key_metrics_table":
            fund_name = data_source.get("fund")
            fund = self.funds[fund_name]
            benchmark_name = data_source.get("benchmark_fund")
            benchmark = self.funds.get(benchmark_name) if benchmark_name else None
            resolved_args = {}
            for k, v in args.items():
                resolved_args[k] = self.resolve(v) if isinstance(v, str) else v
            em = resolved_args.get("end_month")
            if isinstance(em, datetime):
                resolved_args["end_month"] = em.strftime("%Y-%m")
            resolved_args["save"] = False
            resolved_args["benchmark_fund"] = benchmark
            fig = fund.export_key_metrics_table(**resolved_args)
            # matplotlib figure — save to buffer
            buf = io.BytesIO()
            fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0, dpi=200)
            plt.close("all")
            buf.seek(0)
            return buf.read()

        elif method == "export_monthly_table":
            fund_name = data_source.get("fund")
            fund = self.funds[fund_name]
            resolved_args = {}
            for k, v in args.items():
                resolved_args[k] = self.resolve(v) if isinstance(v, str) else v
            em = resolved_args.get("end_month")
            if isinstance(em, datetime):
                resolved_args["end_month"] = em.strftime("%Y-%m")
            resolved_args["save"] = False
            result = fund.export_monthly_table(**resolved_args)
            # result is matplotlib.pyplot — save current figure
            buf = io.BytesIO()
            result.savefig(buf, format="png", bbox_inches="tight", pad_inches=0, dpi=200)
            plt.close("all")
            buf.seek(0)
            return buf.read()

        elif method == "manual_image":
            # Skip manual images
            return None

        return None

    # ------------------------------------------------------------------ #
    #  DOCX operations
    # ------------------------------------------------------------------ #

    def apply_docx_table_cell(self, doc: Document, element: dict):
        """Write a single computed value into a DOCX table cell."""
        loc = element["location"]
        table = doc.tables[loc["table_index"]]
        cell = table.rows[loc["row"]].cells[loc["col"]]
        value = self.compute_value(element["data_source"])
        formatted = self.format_value(value, element.get("format", {}))
        # Preserve formatting: update text of first run in first paragraph
        if cell.paragraphs and cell.paragraphs[0].runs:
            cell.paragraphs[0].runs[0].text = formatted
            # Clear any additional runs
            for run in cell.paragraphs[0].runs[1:]:
                run.text = ""
        else:
            cell.text = formatted

    def apply_docx_text_run(self, doc: Document, element: dict):
        """Replace text runs in a DOCX table cell (for date text)."""
        loc = element["location"]
        table = doc.tables[loc["table_index"]]
        cell = table.rows[loc["row"]].cells[loc["col"]]
        para = cell.paragraphs[0]
        run_start = loc["run_start"]
        run_end = loc["run_end"]
        new_text = self.compute_value(element["data_source"])

        if isinstance(new_text, str):
            # Preserve leading whitespace from the combined original runs
            combined = "".join(
                para.runs[i].text
                for i in range(run_start, min(run_end + 1, len(para.runs)))
            )
            leading_ws = combined[: len(combined) - len(combined.lstrip())]
            # Set the first target run to the new text, clear the rest
            for i in range(run_start, min(run_end + 1, len(para.runs))):
                if i == run_start:
                    para.runs[i].text = (leading_ws or " ") + new_text
                else:
                    para.runs[i].text = ""

    def apply_docx_table_range(self, doc: Document, element: dict):
        """Populate a monthly returns grid in a DOCX table."""
        loc = element["location"]
        ds = element["data_source"]
        fmt = element.get("format", {})

        fund_name = ds.get("fund")
        fund = self.funds[fund_name]
        benchmark_name = self.resolve(ds.get("benchmark_fund"))
        benchmark = self.funds[benchmark_name]
        benchmark_label = self.resolve(ds.get("benchmark_label"))
        args = self.resolve_args(ds.get("args", {}), fund)

        grid = self.build_monthly_returns_grid(
            fund=fund,
            benchmark_fund=benchmark,
            benchmark_label=benchmark_label,
            end_month=args.get("end_month", self.report_month),
            num_years=args.get("num_years"),
            include_ytd=args.get("include_ytd", True),
            include_inception=args.get("include_inception", True),
            fmt=fmt,
        )

        table = doc.tables[loc["table_index"]]
        data_start_row = loc.get("data_start_row", 1)

        # Detect if the grid's top year is newer than the table's existing
        # top year.  When it is, insert row pairs so that old template data
        # shifts down and is preserved beneath the freshly computed rows.
        existing_top_year = None
        if len(table.rows) > data_start_row:
            first_cell_text = table.rows[data_start_row].cells[0].text.strip()
            if first_cell_text.isdigit():
                existing_top_year = int(first_cell_text)

        grid_top_year = None
        if grid and grid[0][0].isdigit():
            grid_top_year = int(grid[0][0])

        if existing_top_year and grid_top_year and grid_top_year > existing_top_year:
            # Each new year adds a fund row + benchmark row
            new_years = grid_top_year - existing_top_year
            self._insert_table_rows(table, data_start_row, new_years * 2)
        else:
            # Insert rows if the grid needs more space than available
            available_data_rows = len(table.rows) - data_start_row
            needed_rows = len(grid)
            if needed_rows > available_data_rows:
                rows_to_add = needed_rows - available_data_rows
                self._insert_table_rows(table, data_start_row, rows_to_add)

        for row_idx, row_data in enumerate(grid):
            table_row_idx = data_start_row + row_idx
            if table_row_idx >= len(table.rows):
                break
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= len(table.columns):
                    break
                cell = table.rows[table_row_idx].cells[col_idx]
                if cell.paragraphs and cell.paragraphs[0].runs:
                    cell.paragraphs[0].runs[0].text = cell_value
                    for run in cell.paragraphs[0].runs[1:]:
                        run.text = ""
                else:
                    cell.text = cell_value

    @staticmethod
    def _insert_table_rows(table, position: int, count: int):
        """Insert *count* new rows into a DOCX table at *position*.

        Copies the XML structure and formatting from the row at *position*
        (the first data row) so new rows inherit cell widths and styles.
        Existing rows from *position* onward are pushed down.
        """
        import copy

        ref_row_elem = table.rows[position]._tr
        parent = ref_row_elem.getparent()

        for _ in range(count):
            new_tr = copy.deepcopy(ref_row_elem)
            # Clear text content in cloned cells
            for tc in new_tr.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                tc.text = ""
            parent.insert(list(parent).index(ref_row_elem), new_tr)

    def apply_docx_image(self, doc: Document, element: dict):
        """Replace an inline image in a DOCX document."""
        from docx.oxml.ns import qn

        image_bytes = self.generate_chart_image(element["data_source"])
        if image_bytes is None:
            return

        dims = element.get("dimensions", {})
        width = dims.get("width_emu")
        height = dims.get("height_emu")

        body_idx = element["location"]["body_element_index"]
        body = doc.element.body
        para_element = body[body_idx]

        # Find existing inline drawings and remove them
        for drawing in para_element.findall(".//" + qn("wp:inline")):
            drawing.getparent().remove(drawing)

        # Find the paragraph object
        from docx.text.paragraph import Paragraph
        para = Paragraph(para_element, doc)

        # Add the new image
        run = para.add_run()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(image_bytes)
            tmp_path = tmp.name

        run.add_picture(tmp_path, width=Emu(width) if width else None, height=Emu(height) if height else None)
        Path(tmp_path).unlink(missing_ok=True)

        # Remove specified body elements (e.g. spacer paragraphs) in reverse
        # order so indices remain valid during removal.
        to_remove = element["location"].get("remove_body_elements", [])
        for idx in sorted(to_remove, reverse=True):
            body.remove(body[idx])

    # ------------------------------------------------------------------ #
    #  PPTX operations
    # ------------------------------------------------------------------ #

    def _find_shape_by_name(self, slide, name: str):
        """Find a shape on a slide by its name."""
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None

    def apply_pptx_slide_text(self, prs: Presentation, element: dict):
        """Replace text in a PPTX shape."""
        loc = element["location"]
        slide = prs.slides[loc["slide_index"]]
        shape = self._find_shape_by_name(slide, loc["shape_name"])
        if shape is None or not shape.has_text_frame:
            print(f"Warning: shape {loc['shape_name']} not found on slide {loc['slide_index']}")
            return

        ds = element["data_source"]
        fmt = element.get("format", {})

        # Compute the replacement value
        if ds.get("method") in ("format_report_date",):
            new_text = self._format_report_date(ds)
        else:
            value = self.compute_value(ds)
            new_text = self.format_value(value, fmt)

        match_pattern = loc.get("match_pattern")
        para_idx = loc.get("paragraph_index")

        if para_idx is not None:
            # Replace text in a specific paragraph
            para = shape.text_frame.paragraphs[para_idx]
            if match_pattern:
                self._replace_para_text(para, match_pattern, new_text)
            else:
                self._set_para_text(para, new_text)
        elif match_pattern:
            # Search all paragraphs for the match pattern
            for para in shape.text_frame.paragraphs:
                if match_pattern in para.text:
                    self._replace_para_text(para, match_pattern, new_text)
                    break

    def _replace_para_text(self, para, pattern: str, replacement: str):
        """Replace a pattern in a paragraph's text, preserving first run formatting."""
        full_text = para.text
        if pattern not in full_text:
            return
        new_full_text = full_text.replace(pattern, replacement)
        # Set the full text on the first run, clear the rest
        if para.runs:
            para.runs[0].text = new_full_text
            for run in para.runs[1:]:
                run.text = ""

    def _set_para_text(self, para, text: str):
        """Set paragraph text, preserving first run formatting."""
        if para.runs:
            para.runs[0].text = text
            for run in para.runs[1:]:
                run.text = ""

    def apply_pptx_slide_table(self, prs: Presentation, element: dict):
        """Populate a PPTX table with computed data."""
        loc = element["location"]
        slide = prs.slides[loc["slide_index"]]
        shape = self._find_shape_by_name(slide, loc["shape_name"])
        if shape is None or not shape.has_table:
            print(f"Warning: table shape {loc['shape_name']} not found on slide {loc['slide_index']}")
            return

        ds = element["data_source"]
        table = shape.table

        if ds.get("method") == "calendar_year_table":
            funds_map = ds.get("funds_map")
            if isinstance(funds_map, str):
                funds_map = self.resolve(funds_map)
            args = self.resolve_args(ds.get("args", {}))
            fmt = element.get("format", {})
            em = args.get("end_month", self.report_month)
            grid = self.build_calendar_year_table(
                funds_map=funds_map,
                start_year=args.get("start_year", 2018),
                end_month=em,
                end_year_offset=args.get("end_year_offset", -1),
                fmt=fmt,
            )
            # Fill table starting at row 1 (row 0 is header)
            for row_idx, row_data in enumerate(grid):
                table_row = row_idx + 1  # skip header
                if table_row >= len(table.rows):
                    break
                for col_idx, val in enumerate(row_data):
                    if col_idx >= len(table.columns):
                        break
                    cell = table.rows[table_row].cells[col_idx]
                    self._set_cell_text(cell, val)

        elif ds.get("method") == "metrics_table":
            fund_name = ds.get("fund")
            fund = self.funds[fund_name]
            args = self.resolve_args(ds.get("args", {}), fund)
            for row_spec in ds.get("rows", []):
                method = row_spec["method"]
                row = row_spec["row"]
                col = row_spec["col"]
                row_fmt = row_spec.get("format", {})

                sub_ds = {
                    "method": method,
                    "fund": fund_name,
                    "args": {
                        "start_month": args.get("start_month"),
                        "end_month": args.get("end_month"),
                    },
                }
                if "benchmark_fund" in row_spec:
                    sub_ds["benchmark_fund"] = row_spec["benchmark_fund"]

                value = self.compute_value(sub_ds)
                formatted = self.format_value(value, row_fmt)
                cell = table.rows[row].cells[col]
                self._set_cell_text(cell, formatted)

    def _set_cell_text(self, cell, text: str):
        """Set text in a PPTX table cell preserving formatting."""
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = text
            for run in cell.text_frame.paragraphs[0].runs[1:]:
                run.text = ""
        else:
            cell.text = text

    def apply_pptx_slide_image(self, prs: Presentation, element: dict):
        """Replace a PICTURE shape in a PPTX slide with a new image."""
        loc = element["location"]
        slide = prs.slides[loc["slide_index"]]
        shape = self._find_shape_by_name(slide, loc["shape_name"])
        if shape is None:
            print(f"Warning: shape {loc['shape_name']} not found on slide {loc['slide_index']}")
            return

        image_bytes = self.generate_chart_image(element["data_source"])
        if image_bytes is None:
            return

        # Replace the image blob
        try:
            shape.image._blob = image_bytes
        except Exception:
            # If direct blob replacement fails, remove and re-add
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(image_bytes)
                tmp_path = tmp.name
            slide.shapes.add_picture(tmp_path, left, top, width, height)
            Path(tmp_path).unlink(missing_ok=True)

    def apply_pptx_cleanup(self, prs: Presentation, element: dict):
        """Remove placeholder overlay text boxes matching a pattern."""
        loc = element["location"]
        pattern = loc.get("match_pattern", r"\[INPUT")
        slides = loc.get("slides", [])

        for slide_idx in slides:
            slide = prs.slides[slide_idx]
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if re.match(pattern, text):
                        shapes_to_remove.append(shape)

            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)


# ====================================================================== #
#  Public API
# ====================================================================== #


def generate_factsheet(
    funds: Dict[str, Fund],
    template_path: str = None,
    spec_path: str = None,
    report_month: str = "2025-12",
    output_path: str = None,
) -> Path:
    """Generate a populated DOCX factsheet from template + spec.

    Parameters
    ----------
    funds : dict
        Dictionary of Fund objects keyed by name.
    template_path : str, optional
        Path to the DOCX template file. Defaults to TEMPLATE_DIR / spec's template_file.
    spec_path : str, optional
        Path to the JSON spec file. Defaults to TEMPLATE_DIR / "docx_spec_en.json".
    report_month : str
        Report month in 'YYYY-MM' format.
    output_path : str, optional
        Output file path. Defaults to output/{fund}_{month}_{lang}.docx.

    Returns
    -------
    Path
        Path to the generated document.
    """
    if spec_path is None:
        spec_path = str(TEMPLATE_DIR / "docx_spec_en.json")

    with open(spec_path, "r") as f:
        spec = json.load(f)

    if template_path is None:
        template_path = str(TEMPLATE_DIR / spec["template_file"])

    doc = Document(template_path)
    engine = TemplateEngine(funds, report_month, spec)

    for element in spec["elements"]:
        elem_type = element["type"]
        try:
            if elem_type == "table_cell":
                engine.apply_docx_table_cell(doc, element)
            elif elem_type == "text_run":
                engine.apply_docx_text_run(doc, element)
            elif elem_type == "table_range":
                engine.apply_docx_table_range(doc, element)
            elif elem_type == "image":
                engine.apply_docx_image(doc, element)
            else:
                print(f"Warning: unknown element type '{elem_type}' for DOCX")
        except Exception as e:
            print(f"Error processing element '{element['id']}': {e}")

    if output_path is None:
        lang = spec["parameters"].get("language", "en")
        fund_name = spec["parameters"]["fund_name"]
        output_path = save_dir / f"{fund_name}_Factsheet_{report_month}_{lang}.docx"

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    print(f"Factsheet saved to {output_path}")
    return output_path


def generate_presentation(
    funds: Dict[str, Fund],
    template_path: str = None,
    spec_path: str = None,
    report_month: str = "2025-12",
    output_path: str = None,
) -> Path:
    """Generate a populated PPTX presentation from template + spec.

    Parameters
    ----------
    funds : dict
        Dictionary of Fund objects keyed by name.
    template_path : str, optional
        Path to the PPTX template file. Defaults to TEMPLATE_DIR / spec's template_file.
    spec_path : str, optional
        Path to the JSON spec file. Defaults to TEMPLATE_DIR / "pptx_spec.json".
    report_month : str
        Report month in 'YYYY-MM' format.
    output_path : str, optional
        Output file path. Defaults to output/{fund}_Presentation_{month}.pptx.

    Returns
    -------
    Path
        Path to the generated presentation.
    """
    if spec_path is None:
        spec_path = str(TEMPLATE_DIR / "pptx_spec.json")

    with open(spec_path, "r") as f:
        spec = json.load(f)

    if template_path is None:
        template_path = str(TEMPLATE_DIR / spec["template_file"])

    prs = Presentation(template_path)
    engine = TemplateEngine(funds, report_month, spec)

    for element in spec["elements"]:
        elem_type = element["type"]
        try:
            if elem_type == "slide_text":
                engine.apply_pptx_slide_text(prs, element)
            elif elem_type == "slide_image":
                engine.apply_pptx_slide_image(prs, element)
            elif elem_type == "slide_table":
                engine.apply_pptx_slide_table(prs, element)
            elif elem_type == "cleanup":
                engine.apply_pptx_cleanup(prs, element)
            else:
                print(f"Warning: unknown element type '{elem_type}' for PPTX")
        except Exception as e:
            print(f"Error processing element '{element['id']}': {e}")

    if output_path is None:
        fund_name = spec["parameters"]["fund_name"]
        output_path = save_dir / f"{fund_name}_Presentation_{report_month}.pptx"

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Presentation saved to {output_path}")
    return output_path
